from pptx import Presentation
from pptx.util import Inches, Pt

class PPT:
    def __init__(self, title):
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(16)
        self.presentation.slide_height = Inches(9)
        self.title = title

    def add_slide(self, title, content):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        slide.shapes.title.text = title
        left = Inches(1)
        top = Inches(2)
        width = Inches(14)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(18)

    def add_results(self, team_results):
        team_data = {}
        for result in team_results:
            team_name = result['team']
            if team_name not in team_data:
                team_data[team_name] = []
            team_data[team_name].append(result)

        for team_name, results in team_data.items():
            title = f"{team_name}"
            content = ""
            for result in results:
                content += (
                    f"종목: {result['stock']}\n"
                    f"수량: {result['quantity']}\n"
                    f"매입단가: {result['purchase_price']}\n"
                    f"평가금액: {result['current_value']}\n"
                    f"평가손익: {result['profit_loss']}\n"
                    f"수익률: {round(result['profit_loss_percentage'],2)}%\n\n"
                )
            content += f"잔액: {results[0]['cash']}\n"
            self.add_slide(title, content)

    def save(self):
        self.presentation.save(f'../ppt/{self.title}.pptx')